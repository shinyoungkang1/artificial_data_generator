"""Artifact writers for multiple data/file formats."""

from __future__ import annotations

import csv
import json
from pathlib import Path
from typing import Any

from .utils import optional_import


def write_json_rows(rows: list[dict[str, Any]], path: str) -> str:
    output = Path(path)
    output.parent.mkdir(parents=True, exist_ok=True)
    payload = {"rows": rows, "row_count": len(rows)}
    output.write_text(json.dumps(payload, indent=2, default=str), encoding="utf-8")
    return str(output)


def write_csv_rows(rows: list[dict[str, Any]], path: str) -> str:
    output = Path(path)
    output.parent.mkdir(parents=True, exist_ok=True)
    if not rows:
        output.write_text("", encoding="utf-8")
        return str(output)

    headers = list(rows[0].keys())
    with output.open("w", newline="", encoding="utf-8") as handle:
        writer = csv.DictWriter(handle, fieldnames=headers)
        writer.writeheader()
        writer.writerows(rows)
    return str(output)


def write_xlsx_rows(rows: list[dict[str, Any]], path: str) -> tuple[str | None, str | None]:
    openpyxl, error = optional_import("openpyxl")
    if openpyxl is None:
        return None, f"openpyxl_not_available:{error}"

    output = Path(path)
    output.parent.mkdir(parents=True, exist_ok=True)

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "raw_data"

    if rows:
        headers = list(rows[0].keys())
        if len(headers) >= 4:
            split = max(2, len(headers) // 2)
            ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=split)
            ws.merge_cells(
                start_row=1,
                start_column=split + 1,
                end_row=1,
                end_column=len(headers),
            )
            ws.cell(row=1, column=1, value="Entity Details")
            ws.cell(row=1, column=split + 1, value="Transaction Details")
        else:
            ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=len(headers))
            ws.cell(row=1, column=1, value="Entity and Transaction Details")

        shifted = 1
        for col_idx, name in enumerate(headers, start=shifted):
            if col_idx % 5 == 0:
                ws.cell(row=2, column=col_idx, value=name.upper())
            else:
                ws.cell(row=2, column=col_idx, value=name)

        row_idx = 3
        for row in rows:
            if row_idx % 23 == 0:
                row_idx += 1
            for col_idx, name in enumerate(headers, start=shifted):
                value = row.get(name)
                ws.cell(row=row_idx, column=col_idx, value=value)
            row_idx += 1

    ws.freeze_panes = "A3"
    wb.save(output)
    return str(output), None


def write_pdf_rows(rows: list[dict[str, Any]], path: str) -> tuple[str | None, str | None]:
    reportlab, error = optional_import("reportlab.pdfgen")
    if reportlab is None:
        return None, f"reportlab_not_available:{error}"

    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas

    output = Path(path)
    output.parent.mkdir(parents=True, exist_ok=True)
    c = canvas.Canvas(str(output), pagesize=letter)
    width, height = letter

    c.setFont("Helvetica-Bold", 12)
    c.drawString(72, height - 60, "Synthetic Company Transactions")
    c.setFont("Helvetica", 9)

    y = height - 90
    for row in rows[:160]:
        text = " | ".join(
            [
                str(row.get("company_id", "")),
                str(row.get("company_name", ""))[:25],
                str(row.get("invoice_id", "")),
                str(row.get("invoice_amount_usd", "")),
                str(row.get("status", "")),
            ]
        )
        c.drawString(72, y, text)
        y -= 11
        if y < 70:
            c.showPage()
            c.setFont("Helvetica", 9)
            y = height - 70

    c.save()
    return str(output), None


def write_pptx_rows(rows: list[dict[str, Any]], path: str) -> tuple[str | None, str | None]:
    pptx, error = optional_import("pptx")
    if pptx is None:
        return None, f"python_pptx_not_available:{error}"

    from pptx import Presentation

    output = Path(path)
    output.parent.mkdir(parents=True, exist_ok=True)

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Synthetic Ops Snapshot"

    sample = rows[:12]
    cols = 4
    rows_count = len(sample) + 1
    table_shape = slide.shapes.add_table(rows_count, cols, 40, 100, 620, 300)
    table = table_shape.table
    table.cell(0, 0).text = "Company"
    table.cell(0, 1).text = "Invoice"
    table.cell(0, 2).text = "Amount"
    table.cell(0, 3).text = "Status"
    for idx, row in enumerate(sample, start=1):
        table.cell(idx, 0).text = str(row.get("company_name", ""))[:25]
        table.cell(idx, 1).text = str(row.get("invoice_id", ""))
        table.cell(idx, 2).text = str(row.get("invoice_amount_usd", ""))
        table.cell(idx, 3).text = str(row.get("status", ""))

    presentation.save(str(output))
    return str(output), None


def render_rows_to_png(rows: list[dict[str, Any]], path: str) -> tuple[str | None, str | None]:
    pil, error = optional_import("PIL")
    if pil is None:
        return None, f"pillow_not_available:{error}"

    from PIL import Image, ImageDraw, ImageFont

    output = Path(path)
    output.parent.mkdir(parents=True, exist_ok=True)
    width = 1800
    height = 2400
    image = Image.new("RGB", (width, height), color=(255, 255, 255))
    draw = ImageDraw.Draw(image)
    font = ImageFont.load_default()

    draw.text((40, 40), "Synthetic Company Transactions (for OCR testing)", fill=(0, 0, 0), font=font)
    y = 80
    for row in rows[:220]:
        line = " | ".join(
            [
                str(row.get("company_name", ""))[:28],
                str(row.get("invoice_id", "")),
                str(row.get("invoice_date", "")),
                str(row.get("invoice_amount_usd", "")),
                str(row.get("status", "")),
            ]
        )
        draw.text((40, y), line, fill=(0, 0, 0), font=font)
        y += 10
        if y > height - 20:
            break

    image.save(output)
    return str(output), None


def pdf_first_page_to_png(pdf_path: str, png_path: str) -> tuple[str | None, str | None]:
    pdf2image, error = optional_import("pdf2image")
    if pdf2image is None:
        return None, f"pdf2image_not_available:{error}"

    try:
        pages = pdf2image.convert_from_path(pdf_path, first_page=1, last_page=1)
    except Exception as exc:
        return None, f"pdf_to_image_failed:{exc}"

    if not pages:
        return None, "pdf_to_image_no_pages"

    output = Path(png_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    pages[0].save(output, format="PNG")
    return str(output), None
